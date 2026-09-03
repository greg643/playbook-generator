#!/usr/bin/env python3
"""
Flag Football Playbook Pipeline
--------------------------------
Takes a .pptx playbook → extracts play images → generates coach cards & wristband PDFs.

Usage:
    python3 playbook_pipeline.py <playbook.pptx> [output_dir]

The script:
1. Reads the PPTX and identifies OFFENSE/DEFENSE sections (or treats a
   headerless deck as offense)
2. Converts slides to high-res images via LibreOffice + pdftoppm
3. Crops each play to its largest rectangle AutoShape
4. Names them 01.png-64.png (offense) and D1.png-D24.png (defense)
5. Feeds them into PlaybookGenerator to create coach cards + wristband PDFs
"""

import os
import sys
import subprocess
import shutil
import json
import tempfile
import re
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Emu

# Import the ink overlay module
from ink_overlay import overlay_ink_on_slides
from input_safety import (
    MAX_DEFENSE_PLAYS,
    MAX_OFFENSE_PLAYS,
    MAX_SLIDES,
    validate_pptx_archive,
    validate_print_play_counts,
)


MAX_SOURCE_IMAGE_DIMENSION = 10000
MAX_SOURCE_IMAGE_PIXELS = 40_000_000
MAX_RENDER_IMAGE_DIMENSION = 1800
# The documented 64-offense + 24-defense capacity must remain possible even
# when every retained crop reaches the per-image render bound.
MAX_TOTAL_SOURCE_IMAGE_PIXELS = (
    (MAX_OFFENSE_PLAYS + MAX_DEFENSE_PLAYS) * MAX_RENDER_IMAGE_DIMENSION ** 2
)
FIELD_RECTANGLE_AUTO_SHAPES = frozenset({
    MSO_SHAPE.RECTANGLE,
    MSO_SHAPE.ROUNDED_RECTANGLE,
    MSO_SHAPE.ROUND_1_RECTANGLE,
    MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
    MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
    MSO_SHAPE.SNIP_1_RECTANGLE,
    MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
    MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
    MSO_SHAPE.SNIP_ROUND_RECTANGLE,
})
OUTPUT_FILENAMES = {
    "offense_coach_card": "offense_coach_card.pdf",
    "offense_wristband": "offense_wristband.pdf",
    "defense_coach_card": "defense_coach_card.pdf",
    "defense_wristband": "defense_wristband.pdf",
}

def wristband_positions(n, column_major=False):
    """Card positions for one wristband cut-out group, chosen by play count so
    partial groups space nicely instead of left-packing: 1-3 = single centered
    row, 4 = 2x2, 5 = 2-1-2 dice, 6 = 3 over 3, 7 = 4 over 3, 8 = classic 4x4
    over two rows. Returns (col, row) in card-pitch units; row 0.5 means
    vertically centered between the two rows. Order is reading order for
    offense (row-major) and column order for defense (column_major=True),
    matching how each side has historically been numbered."""
    n = max(1, min(8, n))
    layouts = {
        1: [(0, 0.5)],
        2: [(0, 0.5), (1, 0.5)],
        3: [(0, 0.5), (1, 0.5), (2, 0.5)],
        4: [(0, 0), (1, 0), (0, 1), (1, 1)],
        5: [(0, 0), (2, 0), (1, 0.5), (0, 1), (2, 1)],
        6: [(0, 0), (1, 0), (2, 0), (0, 1), (1, 1), (2, 1)],
        7: [(0, 0), (1, 0), (2, 0), (3, 0), (0.5, 1), (1.5, 1), (2.5, 1)],
        8: [(0, 0), (1, 0), (2, 0), (3, 0), (0, 1), (1, 1), (2, 1), (3, 1)],
    }
    positions = layouts[n]
    if column_major:
        positions = sorted(positions, key=lambda p: (p[0], p[1]))
    return positions


def wristband_title_allowed(n):
    """A vertical OFFENSE/DEFENSE title fits beside the cards only while the
    arrangement is at most 3 columns wide — with 7-8 cards the grid needs the
    full cut-out width, so the title cannot be drawn."""
    return n <= 6


# Deck-format rule echoed in errors. Separators are optional for offense-only
# decks, and a DEFENSE-only divider can end an implicit opening offense section.
SECTION_HINT = (
    "Every play slide needs a rectangle shape marking the field. A deck with "
    "no OFFENSE separator treats recognizable plays before the first DEFENSE "
    "separator as offense. For explicit sections, start each one with a separator "
    "slide (a slide without a play) whose text "
    "mentions 'Offense' or 'Defense' in any capitalization — '6v6 Offense' works."
)
# ─── STEP 1: Analyze the PPTX ────────────────────────────────────────────────

def analyze_playbook(pptx_path, *, warning_sink=None):
    """
    Walk through slides, detect OFFENSE/DEFENSE sections, identify play slides.
    Returns list of dicts: {slide_index, section, play_number, crop_box_inches}

    The optional warning_sink receives bounded, machine-readable diagnostics
    without changing the long-standing three-value return contract.
    """
    validate_pptx_archive(pptx_path)
    prs = Presentation(pptx_path)
    if len(prs.slides) > MAX_SLIDES:
        raise ValueError(f"PPTX has too many slides (max {MAX_SLIDES})")
    slide_width = prs.slide_width  # EMU
    slide_height = prs.slide_height

    # Pre-scan the complete deck before assigning sections. In particular, a
    # DEFENSE separator does not imply that an earlier OFFENSE separator must
    # exist: many coaching decks simply put offense plays first and add a title
    # card only where defense begins. An explicit OFFENSE separator anywhere in
    # the deck keeps the conservative, explicitly-sectioned behavior.
    slide_records = []
    has_offense_header = False
    has_defense_header = False
    for i, slide in enumerate(prs.slides):
        shapes = list(slide.shapes)
        all_text = " ".join(
            s.text_frame.text for s in shapes if s.has_text_frame
        ).upper()
        crop_box = find_field_rectangle(shapes)
        mentions_offense = "OFFENSE" in all_text
        mentions_defense = "DEFENSE" in all_text
        header_section = None
        if crop_box is None and mentions_offense != mentions_defense:
            header_section = "OFFENSE" if mentions_offense else "DEFENSE"
            if header_section == "OFFENSE":
                has_offense_header = True
            else:
                has_defense_header = True
        slide_records.append({
            "index": i,
            "shapes": shapes,
            "all_text": all_text,
            "crop_box": crop_box,
            "header_section": header_section,
        })

    has_section_headers = has_offense_header or has_defense_header
    plays = []
    current_section = None if has_offense_header else "OFFENSE"
    offense_count = 0
    defense_count = 0
    assumed_offense_count = 0
    skipped_before_section_count = 0
    skipped_no_field_count = 0
    seen_defense_header = False
    seen_section_header = False

    if not has_section_headers:
        print("  No section headers found; treating recognizable play slides as OFFENSE")
    elif not has_offense_header:
        print(
            "  No OFFENSE header found; treating recognizable play slides before "
            "the first DEFENSE header as OFFENSE"
        )

    for record in slide_records:
        i = record["index"]
        shapes = record["shapes"]
        shape_count = len(shapes)
        all_text = record["all_text"]
        crop_box = record["crop_box"]

        # Section separators mention OFFENSE or DEFENSE anywhere in their text,
        # in any capitalization ("6v6 Offense", "OFFENSE PLAYS"). A slide with a
        # field rectangle is never a separator, so a play that merely mentions a
        # section stays a play; a slide mentioning both words is ambiguous and
        # is not a separator either. A repeated section later in the deck
        # continues numbering rather than overwriting 01.png/D1.png from an
        # earlier section.
        if record["header_section"] is not None:
            current_section = record["header_section"]
            seen_section_header = True
            if current_section == "DEFENSE":
                seen_defense_header = True
            print(f"  Slide {i+1}: Section header → {current_section}")
            continue

        skip_keywords = ["PRINT IMAGES", "APPENDIX", "TEMPLATE"]
        is_special_slide = any(kw in all_text for kw in skip_keywords)

        if current_section is None:
            # A field-bearing slide here is likely a play that an explicitly
            # sectioned deck forgot to put after its first divider. Surface the
            # omission after a successful conversion instead of losing it only
            # in the Actions log. Known template/appendix slides stay quiet.
            if crop_box is not None and not is_special_slide:
                skipped_before_section_count += 1
            print(f"  Slide {i+1}: Skipping (before any section, {shape_count} shapes)")
            continue

        if is_special_slide:
            print(f"  Slide {i+1}: Skipping (special slide: {all_text[:40]})")
            continue

        # This is a possible play slide — require an explicit field rectangle.
        # Guessing from the largest arbitrary shape can turn logos/photos into
        # play crops and makes malformed decks appear successful.
        if crop_box is None:
            # Avoid warning on cover/notes that precede the first recognizable
            # play in an implicit opening offense section. After a play or an
            # explicit divider, a no-field slide may be an accidentally omitted
            # play; the UI explains that intentional notes can be ignored.
            if seen_section_header or plays:
                skipped_no_field_count += 1
            print(f"  Slide {i+1}: Skipping (no field rectangle found)")
            continue

        if has_defense_header and not has_offense_header and not seen_defense_header:
            assumed_offense_count += 1

        # Get play name and number from text boxes in the header area.
        # Strategy: find text boxes that START near the field top (within 5%)
        # as the header row. Ignore instruction text lower down.
        play_name = ""
        play_id = ""
        header_bottom_emu = crop_box[1]  # default to field top (no header)
        field_top = crop_box[1]
        field_height = crop_box[3] - crop_box[1]
        header_start_zone = field_top + int(field_height * 0.05)  # must start within top 5%

        header_texts = []
        for s in shapes:
            # Text box object names are localized and user-editable. Rely on
            # the actual text-frame capability while preserving the existing
            # top-of-field positional rule.
            if s.has_text_frame:
                text = s.text_frame.text.strip()
                if text and s.top <= header_start_zone:
                    header_texts.append((s, text))
                    tb = s.top + s.height
                    if tb > header_bottom_emu:
                        header_bottom_emu = tb

        # Separate into play_id (short/numeric) and play_name (longer text)
        for s, text in header_texts:
            if len(text) <= 3 and any(c.isdigit() for c in text):
                play_id = text
            elif len(text) <= 2 and text.isalpha():
                play_id = text  # Defense: A, B, C, D
            else:
                play_name = text

        if current_section == "OFFENSE":
            offense_count += 1
            play_num = offense_count
            filename = f"{play_num:02d}.png"
        else:
            defense_count += 1
            play_num = defense_count
            filename = f"D{play_num}.png"

        # Build label: "1 - Triple Cross" or just "D1"
        if play_id and play_name:
            label = f"{play_id} - {play_name}"
        elif play_id:
            label = play_id
        elif play_name:
            label = play_name
        else:
            label = filename.replace(".png", "")

        plays.append({
            "slide_index": i,
            "section": current_section,
            "play_number": play_num,
            "play_id": play_id,
            "play_name": play_name,
            "label": label,
            "filename": filename,
            "crop_box_emu": crop_box,
            "header_bottom_emu": header_bottom_emu,
        })
        print(f"  Slide {i+1}: {current_section} #{play_num} → {filename} ({play_id} {play_name})")

    if not plays:
        if has_section_headers:
            raise ValueError(
                "Sections were found, but no play slides with a field rectangle "
                "were detected after them. " + SECTION_HINT
            )
        raise ValueError(
            "No play slides with a field rectangle were detected. " + SECTION_HINT
        )

    if warning_sink is not None and assumed_offense_count:
        warning_sink.append({
            "code": "assumed_offense_before_defense",
            "playCount": assumed_offense_count,
        })
    if warning_sink is not None and skipped_before_section_count:
        warning_sink.append({
            "code": "skipped_before_first_divider",
            "slideCount": skipped_before_section_count,
        })
    if warning_sink is not None and skipped_no_field_count:
        warning_sink.append({
            "code": "skipped_no_field_rectangle",
            "slideCount": skipped_no_field_count,
        })

    return plays, slide_width, slide_height


def find_field_rectangle(shapes):
    """
    Find the main field rectangle on a play slide.
    Uses the LARGEST rectangle, which is the full field area.
    On end zone slides, "Rectangle 1" is the yellow end zone (too small),
    while "Rectangle 2" is the full field outline.
    Returns (left, top, right, bottom) in EMU, or None.
    """
    # Prefer the actual PowerPoint rectangle-family geometry. Object names are
    # localized, editable, and not semantic: a genuine rectangle may be renamed
    # "Field", while an oval can be misleadingly named "Rectangle". Keep a
    # conservative name fallback for rectangle AutoShapes whose preset geometry
    # is unknown to python-pptx (for example, a future PowerPoint variant).
    # Group-recursive geometry remains intentionally unsupported because group
    # coordinate transforms require separate crop-box handling.
    rectangles = []
    named_rectangle_fallbacks = []
    for s in shapes:
        try:
            auto_shape_type = s.auto_shape_type
        except (AttributeError, ValueError):
            auto_shape_type = None

        if not (s.width and s.height):
            continue

        area = s.width * s.height
        if auto_shape_type in FIELD_RECTANGLE_AUTO_SHAPES:
            rectangles.append((area, s))
            continue

        # Do not let a known non-rectangle geometry through just because a user
        # renamed it. This fallback is only for an otherwise-unrecognized
        # AutoShape whose PowerPoint-generated name still says "Rectangle".
        if auto_shape_type is None:
            try:
                is_auto_shape = s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
            except (AttributeError, ValueError):
                is_auto_shape = False
            if is_auto_shape and "rectangle" in (getattr(s, "name", "") or "").lower():
                named_rectangle_fallbacks.append((area, s))

    # Only fall back to names when the slide has no geometric rectangle-family
    # match. Within either set, the largest shape is the full field area.
    candidates = rectangles or named_rectangle_fallbacks
    candidates.sort(key=lambda x: x[0], reverse=True)

    if candidates:
        s = candidates[0][1]
        return (s.left, s.top, s.left + s.width, s.top + s.height)

    return None


# ─── STEP 2: Convert PPTX to slide images ────────────────────────────────────

def convert_pptx_to_images(pptx_path, work_dir, dpi=200):
    """
    Convert PPTX → PDF → individual slide PNGs using LibreOffice + pdftoppm.
    Returns path to directory containing slide-NN.png files.
    """
    slides_dir = Path(work_dir) / "slides"
    slides_dir.mkdir(exist_ok=True)

    pdf_path = Path(work_dir) / "playbook.pdf"

    # Step 2a: PPTX → PDF via LibreOffice
    print("\n Converting PPTX → PDF via LibreOffice...")
    pptx_abs = str(Path(pptx_path).resolve())
    pptx_stem = Path(pptx_path).stem

    # Use system LibreOffice (soffice or libreoffice)
    soffice_cmd = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice_cmd:
        # macOS: LibreOffice.app doesn't add to PATH by default
        mac_soffice = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        if os.path.exists(mac_soffice):
            soffice_cmd = mac_soffice
        else:
            raise FileNotFoundError("LibreOffice not found. Install it: brew install --cask libreoffice")
    subprocess.run([
        soffice_cmd, "--headless", "--convert-to", "pdf",
        "--outdir", str(Path(work_dir).resolve()), pptx_abs
    ], check=True, capture_output=True)

    # Find the generated PDF (could be in work_dir or next to input)
    generated_pdf = None
    for candidate in [
        Path(work_dir) / f"{pptx_stem}.pdf",
        Path(pptx_abs).with_suffix(".pdf"),
    ]:
        if candidate.exists():
            generated_pdf = candidate
            break

    if generated_pdf and str(generated_pdf) != str(pdf_path):
        shutil.move(str(generated_pdf), str(pdf_path))
    elif not pdf_path.exists():
        raise FileNotFoundError(f"LibreOffice didn't produce PDF for {pptx_abs}")

    # Step 2b: PDF → PNGs via pdftoppm
    print(f" Converting PDF → PNGs at {dpi} DPI...")
    subprocess.run([
        "pdftoppm", "-png", "-r", str(dpi),
        str(pdf_path),
        str(slides_dir / "slide")
    ], check=True, capture_output=True)

    # List generated files
    slide_images = sorted(slides_dir.glob("slide-*.png"))
    print(f" Generated {len(slide_images)} slide images")
    return slides_dir, slide_images


# ─── STEP 3: Crop play areas from slide images ───────────────────────────────

def crop_plays(plays, slide_images, slide_width_emu, slide_height_emu, output_dir):
    """
    For each play, crop the slide image to the field rectangle (including header).
    """
    output_dir = Path(output_dir)
    output_dir.mkdir(exist_ok=True)

    if not slide_images:
        raise RuntimeError("No slide images were produced from the PowerPoint file")

    # Get image dimensions from first image to compute EMU→pixel ratio
    sample = Image.open(slide_images[0])
    img_w, img_h = sample.size
    sample.close()

    emu_to_px_x = img_w / slide_width_emu
    emu_to_px_y = img_h / slide_height_emu

    normalized_output_names = [
        Path(play["filename"]).name.casefold() for play in plays
    ]
    if len(normalized_output_names) != len(set(normalized_output_names)):
        raise RuntimeError("Detected plays map to duplicate output filenames")

    saved = []
    for play in plays:
        si = play["slide_index"]

        # pdftoppm names slides starting from 01
        slide_num = si + 1  # 1-based
        slide_file = None
        for sf in slide_images:
            name = sf.stem  # "slide-01"
            num_str = name.split("-")[-1]
            if int(num_str) == slide_num:
                slide_file = sf
                break

        if slide_file is None:
            raise RuntimeError(
                f"PowerPoint rendering omitted slide {slide_num}; refusing to create an incomplete playbook"
            )

        img = Image.open(slide_file)

        # Convert crop box from EMU to pixels with small margin for routes
        left_emu, top_emu, right_emu, bottom_emu = play["crop_box_emu"]
        field_w = right_emu - left_emu
        field_h = bottom_emu - top_emu
        margin_x = int(field_w * 0.02)  # 2% horizontal margin
        margin_y = int(field_h * 0.03)  # 3% vertical margin (extra room for routes)

        left_px = max(0, int((left_emu - margin_x) * emu_to_px_x))
        top_px = max(0, int((top_emu - margin_y) * emu_to_px_y))
        right_px = min(img.width, int((right_emu + margin_x) * emu_to_px_x))
        bottom_px = min(img.height, int((bottom_emu + margin_y) * emu_to_px_y))

        cropped = img.crop((left_px, top_px, right_px, bottom_px))

        if cropped.mode != "RGB":
            cropped = cropped.convert("RGB")
        # PDF rendering at 400 DPI can make an ordinary field crop several
        # thousand pixels wide. The printed cards never need that resolution;
        # bounding it here prevents a 64-play deck from predictably exhausting
        # the aggregate generator pixel budget while retaining ample print DPI.
        if max(cropped.size) > MAX_RENDER_IMAGE_DIMENSION:
            cropped.thumbnail(
                (MAX_RENDER_IMAGE_DIMENSION, MAX_RENDER_IMAGE_DIMENSION),
                Image.Resampling.LANCZOS,
            )

        out_path = output_dir / play["filename"]
        cropped.save(out_path, "PNG")
        saved.append(out_path)
        print(f"  {play['filename']:10s} ← Slide {slide_num} ({play['play_id']} {play['play_name']})")

    if len(saved) != len(plays):
        raise RuntimeError("Not every detected play was rendered")
    print(f"\n Saved {len(saved)} play images to {output_dir}/")
    return saved


# ─── STEP 4: PlaybookGenerator (from your existing script) ───────────────────

class PlaybookGenerator:
    """Generates coach cards and wristband PDFs from play images."""

    def __init__(self, images_directory, output_directory="wristband_output"):
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas as pdf_canvas
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader

        self.images_dir = Path(images_directory)
        self.output_dir = Path(output_directory)
        self.output_dir.mkdir(exist_ok=True)

    def fix_image_transparency(self, img):
        if img.mode in ('RGBA', 'LA', 'P'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            if img.mode in ('RGBA', 'LA'):
                background.paste(img, mask=img.split()[-1] if img.mode == 'RGBA' else img.split()[1])
            else:
                background.paste(img)
            return background
        return img.convert('RGB') if img.mode != 'RGB' else img

    def _image_inventory(self):
        """Return canonical play slots and reject files we would otherwise ignore."""
        inventory = {}
        # One-digit offense names may be padded (1.png / 01.png); two-digit
        # names are canonical without an extra leading zero. Bounds are checked
        # separately so the filename grammar cannot drift from input_safety.
        offense_re = re.compile(r"^(0?[1-9]|[1-9][0-9])\.(?:png|jpg)$", re.IGNORECASE)
        defense_re = re.compile(r"^D([1-9]|[1-9][0-9])\.(?:png|jpg)$", re.IGNORECASE)

        for path in sorted(self.images_dir.iterdir()):
            if not path.is_file() or path.suffix.lower() not in {".png", ".jpg"}:
                continue
            offense_match = offense_re.match(path.name)
            defense_match = defense_re.match(path.name)
            if offense_match and int(offense_match.group(1)) <= MAX_OFFENSE_PLAYS:
                number = int(offense_match.group(1))
                slot = ("offense", number)
            elif defense_match and int(defense_match.group(1)) <= MAX_DEFENSE_PLAYS:
                slot = ("defense", int(defense_match.group(1)))
            else:
                raise ValueError(f"Unsupported play image filename: {path.name}")
            if slot in inventory:
                raise ValueError(
                    f"Duplicate play image slot {slot[0]} {slot[1]}: "
                    f"{inventory[slot].name} and {path.name}"
                )
            inventory[slot] = path
        return inventory

    def _load_bounded_image(self, path, total_render_pixels):
        try:
            with Image.open(path) as source:
                if source.format not in {"PNG", "JPEG"}:
                    raise ValueError(f"Unsupported image format for {path.name}")
                width, height = source.size
                pixels = width * height
                if width <= 0 or height <= 0:
                    raise ValueError(f"Invalid image dimensions for {path.name}")
                if width > MAX_SOURCE_IMAGE_DIMENSION or height > MAX_SOURCE_IMAGE_DIMENSION:
                    raise ValueError(
                        f"Play image dimensions are too large for {path.name} "
                        f"(max {MAX_SOURCE_IMAGE_DIMENSION}px per side)"
                    )
                if pixels > MAX_SOURCE_IMAGE_PIXELS:
                    raise ValueError(f"Play image has too many pixels: {path.name}")
                # Images are retained only after being bounded to render size.
                # Account for that retained size rather than the source crop's
                # 400-DPI dimensions, while preserving the per-source cap above
                # and decoding only one source image at a time.
                scale = min(1.0, MAX_RENDER_IMAGE_DIMENSION / max(width, height))
                bounded_width = max(1, round(width * scale))
                bounded_height = max(1, round(height * scale))
                bounded_pixels = bounded_width * bounded_height
                if total_render_pixels + bounded_pixels > MAX_TOTAL_SOURCE_IMAGE_PIXELS:
                    raise ValueError("Combined play images exceed the pixel safety limit")
                source.load()
                image = self.fix_image_transparency(source).copy()
        except (OSError, Image.DecompressionBombError) as exc:
            raise ValueError(f"Could not decode play image: {path.name}") from exc

        if max(image.size) > MAX_RENDER_IMAGE_DIMENSION:
            image.thumbnail(
                (MAX_RENDER_IMAGE_DIMENSION, MAX_RENDER_IMAGE_DIMENSION),
                Image.Resampling.LANCZOS,
            )
        actual_render_pixels = image.width * image.height
        if total_render_pixels + actual_render_pixels > MAX_TOTAL_SOURCE_IMAGE_PIXELS:
            raise ValueError("Combined play images exceed the pixel safety limit")
        return image, total_render_pixels + actual_render_pixels

    def load_images(self):
        offense_images = []
        defense_images = []
        inventory = self._image_inventory()
        total_render_pixels = 0
        for i in range(1, MAX_OFFENSE_PLAYS + 1):
            img_path = inventory.get(("offense", i))
            if img_path:
                img, total_render_pixels = self._load_bounded_image(img_path, total_render_pixels)
                offense_images.append(img)
        for i in range(1, MAX_DEFENSE_PLAYS + 1):
            img_path = inventory.get(("defense", i))
            if img_path:
                img, total_render_pixels = self._load_bounded_image(img_path, total_render_pixels)
                defense_images.append(img)
        return offense_images, defense_images

    def create_coach_card_offense(self, images):
        if not images:
            return
        import io
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader

        pdf_path = self.output_dir / "offense_coach_card.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=landscape(letter))
        page_width, page_height = landscape(letter)
        cols, rows = 4, 4
        margin = 0.5 * inch
        label_space = 0.5 * inch
        grid_width = page_width - 2 * margin - label_space
        grid_height = page_height - 2 * margin
        cell_width = grid_width / cols
        cell_height = grid_height / rows

        page_count = (len(images) + 15) // 16
        for page_num, start_idx in enumerate(range(0, len(images), 16)):
            if page_num > 0:
                c.showPage()

            c.saveState()
            c.setFont("Helvetica-Bold", 24)
            c.translate(margin + label_space/2, page_height/2)
            c.rotate(90)
            c.drawCentredString(0, 0, "OFFENSE")
            c.restoreState()

            page_images = images[start_idx:start_idx + 16]
            page_readers = []
            for img in page_images:
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                page_readers.append(ImageReader(img_buffer))

            for idx, image_reader in enumerate(page_readers):
                row = idx // cols
                col = idx % cols
                x = margin + label_space + col * cell_width
                y = page_height - (margin + (row + 1) * cell_height)
                padding = 3
                c.drawImage(image_reader,
                            x + padding, y + padding,
                            width=cell_width - 2*padding,
                            height=cell_height - 2*padding,
                            preserveAspectRatio=True)
        c.save()
        print(f"  Created: {pdf_path} ({len(images)} plays across {page_count} page(s))")

    def _defense_row_layout(self, n):
        """Return list of plays-per-row for n defense plays: 4→[2,2], 5→[2,1,2], 6→[2,2,2]."""
        if n <= 4:
            return [2, 2]
        elif n == 5:
            return [2, 1, 2]
        else:
            return [2, 2, 2]

    def create_coach_card_defense(self, images):
        if not images:
            return
        import io
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader

        cols = 2

        pdf_path = self.output_dir / "defense_coach_card.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=landscape(letter))
        page_width, page_height = landscape(letter)
        margin = 0.75 * inch
        label_space = 0.5 * inch
        grid_width = page_width - 2 * margin - label_space
        cell_width = grid_width / cols
        grid_height = page_height - 2 * margin

        page_count = (len(images) + 5) // 6
        page_layouts = []
        for page_num, start_idx in enumerate(range(0, len(images), 6)):
            if page_num > 0:
                c.showPage()

            page_images = images[start_idx:start_idx + 6]
            n = len(page_images)
            row_layout = self._defense_row_layout(n)
            page_layouts.append('x'.join(str(r) for r in row_layout))
            num_rows = len(row_layout)
            cell_height = grid_height / num_rows

            c.saveState()
            c.setFont("Helvetica-Bold", 24)
            c.translate(margin + label_space/2, page_height/2)
            c.rotate(90)
            c.drawCentredString(0, 0, "DEFENSE")
            c.restoreState()

            img_idx = 0
            for row_num, count_in_row in enumerate(row_layout):
                for col_num in range(count_in_row):
                    if img_idx >= n:
                        break
                    if count_in_row == 1:
                        x = margin + label_space + (grid_width - cell_width) / 2
                    else:
                        x = margin + label_space + col_num * cell_width
                    y = page_height - (margin + (row_num + 1) * cell_height)
                    img_buffer = io.BytesIO()
                    page_images[img_idx].save(img_buffer, format='PNG')
                    img_buffer.seek(0)
                    padding = 10
                    c.drawImage(ImageReader(img_buffer),
                                x + padding, y + padding,
                                width=cell_width - 2*padding,
                                height=cell_height - 2*padding,
                                preserveAspectRatio=True)
                    img_idx += 1
        c.save()
        print(
            f"  Created: {pdf_path} ({len(images)} plays across {page_count} "
            f"page(s), layouts {', '.join(page_layouts)})"
        )

    # Vertical section-title column on wristband groups (0.25in + 0.05in gap).
    _TITLE_W = 0.25 * 72.0
    _TITLE_GAP = 0.05 * 72.0

    def _draw_group_title(self, c, text, group_x, group_y, group_height):
        c.saveState()
        c.setFont("Helvetica-Bold", 18)
        c.translate(group_x + self._TITLE_W, group_y - group_height / 2)
        c.rotate(90)
        c.drawCentredString(0, 0, text)
        c.restoreState()

    def create_wristband_sheet_offense(self, images, show_title=False):
        if not images:
            return
        import io
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader

        pdf_path = self.output_dir / "offense_wristband.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=landscape(letter))
        page_width, page_height = landscape(letter)

        card_width = 1.0655 * inch
        card_height = 1.0205 * inch
        internal_gap = (3/64) * inch
        group_cols, group_rows = 4, 2
        groups_across, groups_down = 2, 3
        group_width = (group_cols * card_width) + ((group_cols - 1) * internal_gap)
        group_height = (group_rows * card_height) + ((group_rows - 1) * internal_gap)
        group_spacing = 0.5 * inch
        total_width = (groups_across * group_width) + ((groups_across - 1) * group_spacing)
        total_height = (groups_down * group_height) + ((groups_down - 1) * group_spacing)
        start_x = (page_width - total_width) / 2
        start_y = page_height - ((page_height - total_height) / 2)

        # Render ceil(N/8) pages, up to the 64-play offense capacity.
        # Partial pages use the count-adaptive arrangement so a group with 5
        # cards prints as a 2-1-2 dice, 7 as 4-over-3, etc., centered within
        # the same fixed cut-out box as a full group.
        num_pages = (len(images) + 7) // 8
        for page_num in range(num_pages):
            if page_num > 0:
                c.showPage()
            start_idx = page_num * 8
            page_images = images[start_idx:start_idx + 8]
            # Each cut-out group repeats the same page of plays six times. PNG
            # encoding and ImageReader creation are therefore done once per
            # play per page, not once per repeated group.
            page_readers = []
            for img in page_images:
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                page_readers.append(ImageReader(img_buffer))
            positions = wristband_positions(len(page_images))
            grid_width = (
                max(p[0] for p in positions) * (card_width + internal_gap)
                + card_width
            )
            title_ok = show_title and wristband_title_allowed(len(page_images))
            if title_ok:
                cards_area = group_width - self._TITLE_W - self._TITLE_GAP
                cards_x_offset = self._TITLE_W + self._TITLE_GAP + (cards_area - grid_width) / 2
            else:
                cards_x_offset = (group_width - grid_width) / 2
            for group_idx in range(6):
                group_row = group_idx // groups_across
                group_col = group_idx % groups_across
                group_x = start_x + (group_col * (group_width + group_spacing))
                group_y = start_y - (group_row * (group_height + group_spacing))

                # Dashed cutting guide (same size/style as defense wristband)
                c.setStrokeColorRGB(0.3, 0.3, 0.3)
                c.setLineWidth(0.5)
                c.setDash([3, 3])
                c.rect(group_x, group_y - group_height, group_width, group_height)
                c.setDash([])

                if title_ok:
                    self._draw_group_title(c, "OFFENSE", group_x, group_y, group_height)

                for play_idx, (pcol, prow) in enumerate(positions[:len(page_images)]):
                    x = group_x + cards_x_offset + (pcol * (card_width + internal_gap))
                    y = group_y - (prow * (card_height + internal_gap)) - card_height

                    c.drawImage(page_readers[play_idx],
                                x, y, width=card_width, height=card_height,
                                preserveAspectRatio=True, mask='auto')
        c.save()
        print(f"  Created: {pdf_path}")

    def create_wristband_sheet_defense(self, images, show_title=True):
        if not images:
            return
        import io
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.pdfgen import canvas
        from reportlab.lib.units import inch
        from reportlab.lib.utils import ImageReader

        pdf_path = self.output_dir / "defense_wristband.pdf"
        c = canvas.Canvas(str(pdf_path), pagesize=landscape(letter))
        page_width, page_height = landscape(letter)

        # Same card dimensions and gap as offense wristband
        card_width = 1.0655 * inch
        card_height = 1.0205 * inch
        internal_gap = (3/64) * inch

        # Match offense group dimensions (4 cols × 2 rows) for consistent cut-out size
        group_cols_ref = 4
        group_rows_ref = 2
        group_width = (group_cols_ref * card_width) + ((group_cols_ref - 1) * internal_gap)
        group_height = (group_rows_ref * card_height) + ((group_rows_ref - 1) * internal_gap)

        groups_across, groups_down = 2, 3
        group_spacing = 0.5 * inch
        total_width = (groups_across * group_width) + ((groups_across - 1) * group_spacing)
        total_height = (groups_down * group_height) + ((groups_down - 1) * group_spacing)
        start_x = (page_width - total_width) / 2
        start_y = page_height - ((page_height - total_height) / 2)

        num_pages = (len(images) + 7) // 8
        for page_num in range(num_pages):
            if page_num > 0:
                c.showPage()
            start_idx = page_num * 8
            page_images = images[start_idx:start_idx + 8]
            n = len(page_images)

            # Count-adaptive arrangement shared with offense (2-1-2 dice for
            # 5, 3 over 3 for 6, ...), in defense's traditional column order.
            positions = wristband_positions(n, column_major=True)

            # DEFENSE label on the left when enabled and it fits (<= 6 cards);
            # otherwise cards center across the full group width.
            title_ok = show_title and wristband_title_allowed(n)
            defense_grid_width = (
                max(p[0] for p in positions) * (card_width + internal_gap) + card_width
            )
            if title_ok:
                cards_area = group_width - self._TITLE_W - self._TITLE_GAP
                cards_x_offset = self._TITLE_W + self._TITLE_GAP + (cards_area - defense_grid_width) / 2
            else:
                cards_x_offset = (group_width - defense_grid_width) / 2

            image_readers = []
            for img in page_images:
                img_buffer = io.BytesIO()
                img.save(img_buffer, format='PNG')
                img_buffer.seek(0)
                image_readers.append(ImageReader(img_buffer))

            for group_idx in range(6):
                grow = group_idx // groups_across
                gcol = group_idx % groups_across
                group_x = start_x + (gcol * (group_width + group_spacing))
                group_y = start_y - (grow * (group_height + group_spacing))

                # Dashed cutting guide (same size as offense group)
                c.setStrokeColorRGB(0.3, 0.3, 0.3)
                c.setLineWidth(0.5)
                c.setDash([3, 3])
                c.rect(group_x, group_y - group_height, group_width, group_height)
                c.setDash([])

                if title_ok:
                    self._draw_group_title(c, "DEFENSE", group_x, group_y, group_height)

                for img_idx, (pcol, prow) in enumerate(positions[:n]):
                    x = group_x + cards_x_offset + (pcol * (card_width + internal_gap))
                    y = group_y - (prow * (card_height + internal_gap)) - card_height

                    c.drawImage(image_readers[img_idx],
                                x, y, width=card_width, height=card_height,
                                preserveAspectRatio=True, mask='auto')

        c.save()
        print(f"  Created: {pdf_path} ({len(images)} plays across {num_pages} page(s))")

    def generate_all(self, gen_offense=True, gen_defense=True,
                      offense_coach_card=True, offense_wristband=True,
                      defense_coach_card=True, defense_wristband=True,
                      show_offense_title=False, show_defense_title=True):
        expected = set()
        if offense_coach_card:
            expected.add(OUTPUT_FILENAMES["offense_coach_card"])
        if offense_wristband:
            expected.add(OUTPUT_FILENAMES["offense_wristband"])
        if defense_coach_card:
            expected.add(OUTPUT_FILENAMES["defense_coach_card"])
        if defense_wristband:
            expected.add(OUTPUT_FILENAMES["defense_wristband"])
        if not expected:
            raise ValueError("Select at least one output")

        # Known outputs are application-owned. Removing them prevents a prior
        # local run from making a missing output look successful.
        for filename in OUTPUT_FILENAMES.values():
            (self.output_dir / filename).unlink(missing_ok=True)

        print("\nLoading play images...")
        offense_images, defense_images = self.load_images()
        print(f"Found {len(offense_images)} offense plays and {len(defense_images)} defense formations")

        # A one-section deck still produces what it can: outputs requested for
        # a section with no plays are skipped instead of failing the whole job.
        offense_outputs = {OUTPUT_FILENAMES["offense_coach_card"], OUTPUT_FILENAMES["offense_wristband"]}
        defense_outputs = {OUTPUT_FILENAMES["defense_coach_card"], OUTPUT_FILENAMES["defense_wristband"]}
        if not (gen_offense and offense_images):
            for name in sorted(expected & offense_outputs):
                print(f"  No offense plays — skipping {name}")
            expected -= offense_outputs
        if not (gen_defense and defense_images):
            for name in sorted(expected & defense_outputs):
                print(f"  No defense plays — skipping {name}")
            expected -= defense_outputs
        if not expected:
            raise ValueError(
                "No plays were found for any of the selected outputs (for "
                "example, defense outputs selected but the deck has no DEFENSE "
                "section). Adjust the output checkboxes, or fix the deck: "
                + SECTION_HINT
            )

        if gen_offense and offense_images:
            print("\nGenerating offense materials...")
            if offense_coach_card:
                self.create_coach_card_offense(offense_images)
            if offense_wristband:
                self.create_wristband_sheet_offense(offense_images, show_title=show_offense_title)
        if gen_defense and defense_images:
            print("\nGenerating defense materials...")
            if defense_coach_card:
                self.create_coach_card_defense(defense_images)
            if defense_wristband:
                self.create_wristband_sheet_defense(defense_images, show_title=show_defense_title)

        # Only application-owned names count: a pre-existing unrelated PDF in
        # the output directory (CLI runs) must not fail the verification.
        produced = {
            name for name in OUTPUT_FILENAMES.values()
            if (self.output_dir / name).exists()
        }
        if produced != expected:
            missing = sorted(expected - produced)
            unexpected = sorted(produced - expected)
            details = []
            if missing:
                details.append("missing: " + ", ".join(missing))
            if unexpected:
                details.append("unexpected: " + ", ".join(unexpected))
            raise RuntimeError("Generated output set did not match the request (" + "; ".join(details) + ")")

        print(f"\nDone! Output in: {self.output_dir}/")
        for pdf in sorted(self.output_dir.glob("*.pdf")):
            print(f"  {pdf.name}")
        return sorted(produced)


# ─── MAIN ────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 playbook_pipeline.py <playbook.pptx> [output_dir] [--sections offense|defense|both] [--mode standard|screenshot]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 and not sys.argv[2].startswith("--") else "playbook_output"

    # Parse --sections flag (shorthand: offense, defense, or both)
    sections = "both"
    if "--sections" in sys.argv:
        idx = sys.argv.index("--sections")
        if idx + 1 < len(sys.argv):
            sections = sys.argv[idx + 1].lower()

    # Parse --outputs flag (granular: comma-separated list of output types)
    # e.g. --outputs offense_coach_card,offense_wristband,defense_wristband
    # If not specified, derive from --sections (all outputs for selected sections)
    all_outputs = {"offense_coach_card", "offense_wristband", "defense_coach_card", "defense_wristband"}
    if "--outputs" in sys.argv:
        idx = sys.argv.index("--outputs")
        if idx + 1 < len(sys.argv):
            selected_outputs = set(sys.argv[idx + 1].lower().split(","))
        else:
            selected_outputs = all_outputs
    else:
        selected_outputs = set()
        if sections in ("both", "offense"):
            selected_outputs |= {"offense_coach_card", "offense_wristband"}
        if sections in ("both", "defense"):
            selected_outputs |= {"defense_coach_card", "defense_wristband"}

    gen_offense = "offense_coach_card" in selected_outputs or "offense_wristband" in selected_outputs
    gen_defense = "defense_coach_card" in selected_outputs or "defense_wristband" in selected_outputs
    offense_coach_card = "offense_coach_card" in selected_outputs
    offense_wristband = "offense_wristband" in selected_outputs
    defense_coach_card = "defense_coach_card" in selected_outputs
    defense_wristband = "defense_wristband" in selected_outputs

    # Parse --titles flag: which wristbands draw the vertical section title
    # (comma list of offense/defense; "none" for neither; default defense).
    # Titles render only on cut-out groups with 6 or fewer cards.
    titles = {"defense"}
    if "--titles" in sys.argv:
        idx = sys.argv.index("--titles")
        if idx + 1 < len(sys.argv):
            titles = {
                t for t in sys.argv[idx + 1].lower().split(",")
                if t in ("offense", "defense")
            }

    # Parse --mode flag: "standard" (200 DPI + ink overlay) or "screenshot" (600 DPI, no ink overlay)
    mode = "standard"
    if "--mode" in sys.argv:
        idx = sys.argv.index("--mode")
        if idx + 1 < len(sys.argv):
            mode = sys.argv[idx + 1].lower()
    render_dpi = 600 if mode == "screenshot" else 400

    # Each invocation gets an isolated workspace.  Runs are intentionally kept
    # under _playbook_work for local inspection, but no run can reuse another
    # deck's slides, ink overlays, or numbered play images.
    work_root = Path("_playbook_work")
    work_root.mkdir(exist_ok=True)
    work_dir = Path(tempfile.mkdtemp(prefix="run-", dir=work_root))
    plays_dir = work_dir / "plays"
    plays_dir.mkdir(exist_ok=True)

    print(f"{'='*60}")
    print(f"Flag Football Playbook Pipeline")
    print(f"{'='*60}")
    print(f"Input:  {pptx_path}")
    print(f"Output: {output_dir}/")
    print(f"Mode:   {mode} ({render_dpi} DPI{', no ink overlay' if mode == 'screenshot' else ''})")
    print()

    # Step 1: Analyze
    print("STEP 1: Analyzing playbook structure...")
    analysis_warnings = []
    plays, slide_w, slide_h = analyze_playbook(
        pptx_path,
        warning_sink=analysis_warnings,
    )
    validate_print_play_counts(plays)
    offense_plays = [p for p in plays if p["section"] == "OFFENSE"]
    defense_plays = [p for p in plays if p["section"] == "DEFENSE"]
    print(f"\n  Found {len(offense_plays)} offense plays, {len(defense_plays)} defense plays")

    # Step 2: Convert to images (skip if already done)
    slides_dir = work_dir / "slides"
    existing_slides = sorted(slides_dir.glob("slide-*.png")) if slides_dir.exists() else []
    # Only use existing slides that don't have _with_ink suffix (base slides)
    base_slides = [s for s in existing_slides if "_with_ink" not in s.name]
    if base_slides and len(base_slides) >= len(plays):
        print(f"\nSTEP 2: Using {len(base_slides)} existing slide images (skipping conversion)")
        slide_images = base_slides
    else:
        print(f"\nSTEP 2: Converting slides to images at {render_dpi} DPI...")
        slides_dir, slide_images = convert_pptx_to_images(pptx_path, work_dir, dpi=render_dpi)

    # Step 2.5: Overlay ink annotations (standard mode only)
    if mode == "screenshot":
        print("\nSTEP 2.5: Screenshot mode — skipping ink overlay (using LibreOffice native rendering)")
    else:
        ink_files = list(slides_dir.glob("*_with_ink.png")) if slides_dir.exists() else []
        if ink_files:
            print(f"\nSTEP 2.5: Ink overlays already applied ({len(ink_files)} files), skipping...")
        else:
            print("\nSTEP 2.5: Overlaying ink annotations (hand-drawn routes)...")
            ink_output = overlay_ink_on_slides(
                pptx_path=str(Path(pptx_path).resolve()),
                # The overlay implementation reads directly from ZipFile; no
                # extraction is needed (and avoiding it removes a zip-bomb sink).
                pptx_unzipped_path="",
                slides_dir=str(slides_dir),
                approach='B',
                use_fallback_if_failed=True,
                dpi=render_dpi
            )
            # Replace original slide images with ink-overlaid versions
            for slide_num, ink_path in ink_output.items():
                ink_img_path = Path(ink_path)
                for sf in slide_images:
                    num_str = sf.stem.split("-")[-1]
                    if int(num_str) == slide_num:
                        shutil.copy2(str(ink_img_path), str(sf))
                        print(f"  Replaced slide-{slide_num:02d}.png with ink-overlaid version")
                        break
            print(f"  Overlaid ink on {len(ink_output)} slides")

    # Step 3: Crop plays
    print("\nSTEP 3: Cropping play images...")
    crop_plays(plays, slide_images, slide_w, slide_h, plays_dir)

    # Step 4: Generate PDFs
    print(f"\nSTEP 4: Generating coach cards and wristbands (sections: {sections})...")
    generator = PlaybookGenerator(str(plays_dir), output_dir)
    generator.generate_all(gen_offense=gen_offense, gen_defense=gen_defense,
                           offense_coach_card=offense_coach_card, offense_wristband=offense_wristband,
                           defense_coach_card=defense_coach_card, defense_wristband=defense_wristband,
                           show_offense_title="offense" in titles,
                           show_defense_title="defense" in titles)

    # Cleanup
    print(f"\nPlay images saved in: {plays_dir}/")
    print(f"Final PDFs saved in: {output_dir}/")
    return {"warnings": analysis_warnings}


if __name__ == "__main__":
    main()
