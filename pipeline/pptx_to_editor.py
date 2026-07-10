#!/usr/bin/env python3
"""
pptx_to_editor.py — Convert a flag-football PPTX playbook into the play-editor
playbook JSON (schema 1, see dashboard/editor.html normalizePlay()).

Usage:
    python pptx_to_editor.py <playbook.pptx> <out.json>

Per play, the FIELD region in slide EMU is
    x: crop_left..crop_right,  y: header_bottom_emu..crop_bottom
and every shape coordinate is normalized into 0..1 within that region.

Imported per play:
  - CHIPS  : shapes whose text is exactly 1-5/QB (offense) or 1-5/N (defense)
  - ROUTES : hand-drawn InkML strokes (anchored to nearest chip) and native
             line/connector shapes that start at a chip (elbow pieces are
             chained into one polyline first)
  - LINES  : remaining native lines/connectors (line of scrimmage, zone lines)
  - LABELS : text shapes below the header that aren't chip labels
  - BALLS  : football picture shapes not consumed by a route tip
"""

import json
import math
import re
import sys
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).resolve().parent))
from playbook_pipeline import analyze_playbook  # noqa: E402
from ink_overlay import extract_ink_strokes, NS  # noqa: E402

# ── constants ────────────────────────────────────────────────────────────────

PALETTE = {
    '#FF0000': (255, 0, 0),
    '#000000': (0, 0, 0),
    '#1F6E8C': (31, 110, 140),
    '#0070C0': (0, 112, 192),
    '#00B050': (0, 176, 80),
}
OFFENSE_KEYS = {'1', '2', '3', '4', '5', 'QB'}
DEFENSE_KEYS = {'1', '2', '3', '4', '5', 'N'}
ARROW_TYPES = {'triangle', 'arrow', 'stealth', 'diamond', 'oval'}
LINE_GEOMS = {'line', 'straightConnector1', 'bentConnector2', 'bentConnector3',
              'curvedConnector2', 'curvedConnector3'}

CHAIN_TOL_EMU = 60000          # endpoint gap to merge elbow pieces into one polyline
FULL_WIDTH_FRAC = 0.90         # x-extent >= this fraction of field width -> always a LINE
ROUTE_ATTACH_EMU = 500000      # chip-attach radius (~1.4 chip radii); farther = loose LINE
BALL_TIP_FRAC = 0.06           # football within this fraction of field diagonal of a tip
INK_MARKER_DIAG_FRAC = 0.055   # ink stroke smaller than this = arrowhead scribble
INK_MARKER_ATTACH_FRAC = 0.08  # marker attaches to a body endpoint within this
MAX_ROUTE_POINTS = 18

THEME_ALIAS = {'tx1': 'dk1', 'bg1': 'lt1', 'tx2': 'dk2', 'bg2': 'lt2'}


# ── small helpers ────────────────────────────────────────────────────────────

def dist(a, b):
    return math.hypot(a[0] - b[0], a[1] - b[1])


def nearest_palette(rgb):
    best, best_d = '#000000', float('inf')
    for hexc, (r, g, b) in PALETTE.items():
        d = (rgb[0] - r) ** 2 + (rgb[1] - g) ** 2 + (rgb[2] - b) ** 2
        if d < best_d:
            best, best_d = hexc, d
    return best


def hex_to_rgb(h):
    h = h.lstrip('#')
    if len(h) != 6:
        return (0, 0, 0)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def load_theme_colors(pptx_zip):
    """Map theme color names (dk1, accent1, ...) to hex strings."""
    colors = {}
    try:
        root = None
        for name in pptx_zip.namelist():
            if re.match(r'ppt/theme/theme1\.xml$', name):
                import xml.etree.ElementTree as ET
                root = ET.fromstring(pptx_zip.read(name))
                break
        if root is None:
            return colors
        scheme = root.find('.//a:clrScheme', NS)
        if scheme is None:
            return colors
        for child in scheme:
            tag = child.tag.split('}')[-1]
            srgb = child.find('a:srgbClr', NS)
            sysc = child.find('a:sysClr', NS)
            if srgb is not None:
                colors[tag] = srgb.get('val')
            elif sysc is not None:
                colors[tag] = sysc.get('lastClr', '000000')
    except Exception:
        pass
    return colors


def resolve_scheme(name, theme):
    name = THEME_ALIAS.get(name, name)
    return theme.get(name, '000000')


def line_props(el, theme):
    """Extract (color_rgb_or_None, dash, arrow_at_start, arrow_at_end, width_emu)
    from a shape/connector element's a:ln (with p:style lnRef fallback)."""
    color = None
    dash = False
    a_start = a_end = False
    width = None
    ln = el.find('p:spPr/a:ln', NS)
    if ln is not None:
        if ln.get('w'):
            try:
                width = int(ln.get('w'))
            except ValueError:
                pass
        srgb = ln.find('a:solidFill/a:srgbClr', NS)
        scheme = ln.find('a:solidFill/a:schemeClr', NS)
        if srgb is not None:
            color = hex_to_rgb(srgb.get('val', '000000'))
        elif scheme is not None:
            color = hex_to_rgb(resolve_scheme(scheme.get('val', 'dk1'), theme))
        pd = ln.find('a:prstDash', NS)
        if pd is not None and pd.get('val', 'solid') != 'solid':
            dash = True
        if ln.find('a:custDash', NS) is not None:
            dash = True
        he = ln.find('a:headEnd', NS)
        te = ln.find('a:tailEnd', NS)
        if he is not None and he.get('type') in ARROW_TYPES:
            a_start = True
        if te is not None and te.get('type') in ARROW_TYPES:
            a_end = True
    if color is None:
        style = el.find('p:style/a:lnRef/a:schemeClr', NS)
        if style is not None:
            color = hex_to_rgb(resolve_scheme(style.get('val', 'dk1'), theme))
    return color, dash, a_start, a_end, width


def connector_local_points(prst, cx, cy, adj_vals):
    """Approximate a connector's geometry as local points in its (cx, cy) box."""
    if prst in ('line', 'straightConnector1'):
        return [(0.0, 0.0), (float(cx), float(cy))]
    if prst == 'bentConnector2':
        return [(0.0, 0.0), (float(cx), 0.0), (float(cx), float(cy))]
    if prst == 'bentConnector3':
        adj = adj_vals.get('adj1', 50000) / 100000.0
        x2 = adj * cx
        return [(0.0, 0.0), (x2, 0.0), (x2, float(cy)), (float(cx), float(cy))]
    if prst in ('curvedConnector2', 'curvedConnector3'):
        # curvedConnector3: leaves start horizontally, S-curves through a
        # vertical mid-line at x = adj1% of width, arrives horizontally.
        adj = adj_vals.get('adj1', 50000) / 100000.0
        x2 = adj * cx
        pts = []
        # quadratic bezier 1: (0,0) ctrl (x2,0) -> (x2, cy/2)
        for i in range(5):
            t = i / 4.0
            x = (1 - t) ** 2 * 0.0 + 2 * (1 - t) * t * x2 + t ** 2 * x2
            y = 2 * (1 - t) * t * 0.0 + t ** 2 * (cy / 2.0)
            pts.append((x, y))
        # quadratic bezier 2: (x2, cy/2) ctrl (x2, cy) -> (cx, cy)
        for i in range(1, 5):
            t = i / 4.0
            x = (1 - t) ** 2 * x2 + 2 * (1 - t) * t * x2 + t ** 2 * cx
            y = (1 - t) ** 2 * (cy / 2.0) + 2 * (1 - t) * t * cy + t ** 2 * cy
            pts.append((x, y))
        return pts
    return [(0.0, 0.0), (float(cx), float(cy))]


def transform_local_points(pts, cx, cy, flip_h, flip_v, rot_deg, left, top, xf):
    """Flip about the box, rotate about the box center, translate to parent
    coords, then apply the accumulated group transform xf=(ox, oy, sx, sy)."""
    ox, oy, sx, sy = xf
    ctr_x, ctr_y = cx / 2.0, cy / 2.0
    out = []
    cos_t = math.cos(math.radians(rot_deg))
    sin_t = math.sin(math.radians(rot_deg))
    for (x, y) in pts:
        if flip_h:
            x = cx - x
        if flip_v:
            y = cy - y
        dx, dy = x - ctr_x, y - ctr_y
        rx = dx * cos_t - dy * sin_t
        ry = dx * sin_t + dy * cos_t
        px = left + ctr_x + rx
        py = top + ctr_y + ry
        out.append((ox + px * sx, oy + py * sy))
    return out


def _perp_dist(p, a, b):
    """Perpendicular distance from p to segment a-b."""
    dx, dy = b[0] - a[0], b[1] - a[1]
    l2 = dx * dx + dy * dy
    if l2 == 0:
        return dist(p, a)
    t = max(0.0, min(1.0, ((p[0] - a[0]) * dx + (p[1] - a[1]) * dy) / l2))
    return dist(p, (a[0] + t * dx, a[1] + t * dy))


def _rdp(pts, eps):
    """Ramer-Douglas-Peucker: keeps high-curvature points (hooks, curls)
    that uniform arc-length sampling flattens away."""
    if len(pts) < 3:
        return list(pts)
    dmax, idx = 0.0, 0
    for i in range(1, len(pts) - 1):
        d = _perp_dist(pts[i], pts[0], pts[-1])
        if d > dmax:
            dmax, idx = d, i
    if dmax <= eps:
        return [pts[0], pts[-1]]
    left = _rdp(pts[:idx + 1], eps)
    right = _rdp(pts[idx:], eps)
    return left[:-1] + right


def downsample(pts, max_n=MAX_ROUTE_POINTS, eps_hint=None):
    """Curvature-preserving simplification (RDP), escalating tolerance
    until the point budget fits. Keeps both endpoints."""
    if len(pts) <= max_n:
        return list(pts)
    span = max(max(p[0] for p in pts) - min(p[0] for p in pts),
               max(p[1] for p in pts) - min(p[1] for p in pts)) or 1.0
    eps = eps_hint if eps_hint else span * 0.008
    out = _rdp(list(pts), eps)
    while len(out) > max_n:
        eps *= 1.6
        out = _rdp(list(pts), eps)
    return out




def merge_block_caps(routes, lines):
    """A blocking 'T' in the deck arrives as a short stub plus a detached
    perpendicular cap line; fold the cap into the stub as end='block'."""
    def seg_dir(pts):
        (x1, y1), (x2, y2) = pts[0], pts[-1]
        d = math.hypot(x2 - x1, y2 - y1) or 1.0
        return (x2 - x1) / d, (y2 - y1) / d

    remaining = []
    for ln in lines:
        pts = ln['points']
        if len(pts) == 2 and ln.get('end', 'none') == 'none':
            (x1, y1), (x2, y2) = pts
            cap_len = math.hypot(x2 - x1, y2 - y1)
            if cap_len <= 0.09:
                mid = ((x1 + x2) / 2, (y1 + y2) / 2)
                cd = seg_dir(pts)
                matched = False
                for rt in routes:
                    tip = rt['points'][-1]
                    if math.hypot(mid[0] - tip[0], mid[1] - tip[1]) <= 0.035 and rt.get('end') != 'ball':
                        rd = seg_dir(rt['points'][-2:])
                        if abs(rd[0] * cd[0] + rd[1] * cd[1]) < 0.45:  # near-perpendicular
                            rt['end'] = 'block'
                            matched = True
                            break
                if matched:
                    continue
        remaining.append(ln)
    return routes, remaining


# ── shape collection ─────────────────────────────────────────────────────────

def get_xfrm_attrs(el):
    xfrm = el.find('.//a:xfrm', NS)
    if xfrm is None:
        return 0, False, False
    rot = int(xfrm.get('rot', '0') or 0) / 60000.0
    return rot, xfrm.get('flipH') == '1', xfrm.get('flipV') == '1'


def collect_shapes(shapes, xf, out):
    """Recursively flatten shapes; group children carry an accumulated
    (offset, scale) transform mapping their coords into slide EMU."""
    for s in shapes:
        if s.shape_type == MSO_SHAPE_TYPE.GROUP:
            el = s._element
            gx = el.find('p:grpSpPr/a:xfrm', NS)
            if gx is None:
                continue
            off = gx.find('a:off', NS)
            ext = gx.find('a:ext', NS)
            ch_off = gx.find('a:chOff', NS)
            ch_ext = gx.find('a:chExt', NS)
            if off is None or ext is None:
                continue
            g_l, g_t = int(off.get('x', 0)), int(off.get('y', 0))
            g_w, g_h = int(ext.get('cx', 1)) or 1, int(ext.get('cy', 1)) or 1
            c_x = int(ch_off.get('x', 0)) if ch_off is not None else 0
            c_y = int(ch_off.get('y', 0)) if ch_off is not None else 0
            c_w = (int(ch_ext.get('cx', g_w)) if ch_ext is not None else g_w) or g_w
            c_h = (int(ch_ext.get('cy', g_h)) if ch_ext is not None else g_h) or g_h
            ox, oy, sx, sy = xf
            nsx = sx * (g_w / c_w)
            nsy = sy * (g_h / c_h)
            nox = ox + sx * (g_l - c_x * (g_w / c_w))
            noy = oy + sy * (g_t - c_y * (g_h / c_h))
            collect_shapes(s.shapes, (nox, noy, nsx, nsy), out)
        else:
            out.append((s, xf))
    return out


def shape_bbox(s, xf):
    ox, oy, sx, sy = xf
    left = ox + (s.left or 0) * sx
    top = oy + (s.top or 0) * sy
    w = (s.width or 0) * sx
    h = (s.height or 0) * sy
    return left, top, w, h


def shape_text(s):
    try:
        if s.has_text_frame:
            return s.text_frame.text.strip()
    except Exception:
        pass
    return ''


def first_run_style(s):
    """(font_size_emu_or_None, color_hex_or_None) from the first styled run."""
    size = None
    color = None
    try:
        for para in s.text_frame.paragraphs:
            for run in para.runs:
                if size is None and run.font.size is not None:
                    size = int(run.font.size)
                if color is None:
                    try:
                        rgb = run.font.color.rgb
                        if rgb is not None:
                            color = f'#{rgb}'
                    except Exception:
                        pass
                if size is not None and color is not None:
                    return size, color
    except Exception:
        pass
    return size, color


# ── native segment chaining ──────────────────────────────────────────────────

def chain_segments(segments):
    """Merge segments whose endpoints touch (same color+dash) into elbow
    polylines. Never merge through an arrowhead."""
    segs = [dict(s) for s in segments]
    merged = True
    while merged:
        merged = False
        for i in range(len(segs)):
            if merged:
                break
            for j in range(len(segs)):
                if i == j:
                    continue
                a, b = segs[i], segs[j]
                if a['color'] != b['color'] or a['dash'] != b['dash']:
                    continue
                # candidate joins: (a end -> b start) needs a.a1 and b.a0 clear
                if dist(a['pts'][-1], b['pts'][0]) <= CHAIN_TOL_EMU and not a['a1'] and not b['a0']:
                    a['pts'] = a['pts'] + b['pts'][1:]
                    a['a1'] = b['a1']
                elif dist(a['pts'][-1], b['pts'][-1]) <= CHAIN_TOL_EMU and not a['a1'] and not b['a1']:
                    a['pts'] = a['pts'] + list(reversed(b['pts']))[1:]
                    a['a1'] = b['a0']
                elif dist(a['pts'][0], b['pts'][-1]) <= CHAIN_TOL_EMU and not a['a0'] and not b['a1']:
                    a['pts'] = b['pts'] + a['pts'][1:]
                    a['a0'] = b['a0']
                elif dist(a['pts'][0], b['pts'][0]) <= CHAIN_TOL_EMU and not a['a0'] and not b['a0']:
                    a['pts'] = list(reversed(b['pts'])) + a['pts'][1:]
                    a['a0'] = b['a1']
                else:
                    continue
                a['width'] = a['width'] or b['width']
                segs.pop(j)
                merged = True
                break
    return segs


# ── per-play conversion ──────────────────────────────────────────────────────

def convert_play(play_meta, prs, pptx_zip, theme, section):
    slide = prs.slides[play_meta['slide_index']]
    crop_l, crop_t, crop_r, crop_b = play_meta['crop_box_emu']
    field_x0 = crop_l
    field_y0 = play_meta['header_bottom_emu']
    field_w = max(1, crop_r - crop_l)
    field_h = max(1, crop_b - field_y0)
    field_diag = math.hypot(field_w, field_h)

    def norm(pt):
        return [
            round(min(1.0, max(0.0, (pt[0] - field_x0) / field_w)), 4),
            round(min(1.0, max(0.0, (pt[1] - field_y0) / field_h)), 4),
        ]

    chip_keys = DEFENSE_KEYS if section == 'defense' else OFFENSE_KEYS

    flat = collect_shapes(slide.shapes, (0.0, 0.0, 1.0, 1.0), [])

    chips = {}          # key -> (x_emu, y_emu)
    labels = []
    balls = []          # dicts {c:(x,y), used:False}
    segments = []       # native line-ish pieces

    for s, xf in flat:
        el = s._element
        left, top, w, h = shape_bbox(s, xf)
        cx_pt = (left + w / 2.0, top + h / 2.0)
        if cx_pt[1] < field_y0:            # header area: number/name cells etc.
            continue

        geom_el = el.find('.//a:prstGeom', NS)
        geom = geom_el.get('prst') if geom_el is not None else None
        text = shape_text(s)

        if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
            balls.append({'c': cx_pt, 'used': False})
            continue

        is_line_geom = (geom in LINE_GEOMS) or (el.tag == f'{{{NS["p"]}}}cxnSp')
        if is_line_geom:
            color_rgb, dash_flag, a0, a1, w_emu = line_props(el, theme)
            rot, flip_h, flip_v = get_xfrm_attrs(el)
            adj_vals = {}
            if geom_el is not None:
                for gd in geom_el.findall('a:avLst/a:gd', NS):
                    m = re.match(r'val (-?\d+)', gd.get('fmla', ''))
                    if m:
                        adj_vals[gd.get('name')] = int(m.group(1))
            local = connector_local_points(geom or 'line', s.width or 0, s.height or 0, adj_vals)
            pts = transform_local_points(local, s.width or 0, s.height or 0,
                                         flip_h, flip_v, rot,
                                         s.left or 0, s.top or 0, xf)
            if dist(pts[0], pts[-1]) < 1000 and len(pts) == 2:
                continue  # degenerate
            segments.append({
                'pts': pts,
                'color': nearest_palette(color_rgb) if color_rgb else ('#1F6E8C' if dash_flag else '#000000'),
                'dash': dash_flag,
                'a0': a0, 'a1': a1,
                'width': w_emu,
            })
            continue

        if text and text in chip_keys:
            chips[text] = cx_pt
            continue

        if text:
            # Keep line breaks: the editor renders multi-line labels and
            # auto-wraps long ones inside the field.
            clean = re.sub(r'[\r\x0b]+', '\n', text)
            clean = re.sub(r'[ \t]{2,}', ' ', clean).strip()[:160]
            if clean:
                size_emu, color_hex = first_run_style(s)
                frac = (size_emu / field_h) if size_emu else 0.055
                labels.append({
                    'x': norm(cx_pt)[0], 'y': norm(cx_pt)[1],
                    'text': clean,
                    # PowerPoint's default text color is black (tx1), so use
                    # black when no explicit run color is set
                    'color': color_hex or '#000000',
                    'size': round(min(0.2, max(0.02, frac)), 4),
                })
        # rectangles without text (field border, header, pylons) are skipped

    routes = []
    lines = []

    def consume_ball(tip):
        for b in balls:
            if not b['used'] and dist(b['c'], tip) <= BALL_TIP_FRAC * field_diag:
                b['used'] = True
                return True
        return False

    def nearest_chip(pt):
        best_key, best_d = None, float('inf')
        for key, c in chips.items():
            d = dist(pt, c)
            if d < best_d:
                best_key, best_d = key, d
        return best_key, best_d

    # ---- native segments: chain elbows, then classify route vs line --------
    for seg in chain_segments(segments):
        pts = seg['pts']
        x_extent = max(p[0] for p in pts) - min(p[0] for p in pts)
        full_width = x_extent >= FULL_WIDTH_FRAC * field_w

        route_key = None
        if not full_width and chips:
            thr = ROUTE_ATTACH_EMU
            k0, d0 = nearest_chip(pts[0])
            k1, d1 = nearest_chip(pts[-1])
            if min(d0, d1) <= thr:
                if d1 < d0:                 # flip so the chip end comes first
                    pts = list(reversed(pts))
                    seg['a0'], seg['a1'] = seg['a1'], seg['a0']
                    route_key = k1
                else:
                    route_key = k0

        if route_key is not None:
            if seg['a1']:
                end = 'arrow'
            elif consume_ball(pts[-1]):
                end = 'ball'
            else:
                end = 'none'
            routes.append({
                'chip': route_key,
                'color': seg['color'],
                'dash': seg['dash'],
                'end': end,
                'corner': 'sharp',   # native connectors have crisp elbows
                'points': [norm(p) for p in downsample(pts)],
            })
        else:
            if seg['a0'] and not seg['a1']:
                pts = list(reversed(pts))
                seg['a0'], seg['a1'] = seg['a1'], seg['a0']
            line = {
                'color': seg['color'],
                'dash': seg['dash'],
                'end': 'arrow' if seg['a1'] else 'none',
                'corner': 'sharp',
                'points': [norm(p) for p in downsample(pts)],
            }
            if seg['dash']:
                line['width'] = 10
                line['dashArray'] = '28 18'
            # Solid lines take the editor's uniform stroke (13): PowerPoint decks
            # carry slightly-varying point sizes that print as inconsistent weight.
            lines.append(line)

    # ---- ink strokes: bodies become routes, tiny scribbles mark arrowheads --
    strokes = []
    try:
        raw = extract_ink_strokes('', play_meta['slide_index'] + 1, pptx_zip)
    except Exception:
        raw = []
    for stroke_pts, (bx, by, bw, bh), color_hex, _bwidth in raw:
        if len(stroke_pts) < 2:
            continue
        min_x = min(p[0] for p in stroke_pts)
        max_x = max(p[0] for p in stroke_pts)
        min_y = min(p[1] for p in stroke_pts)
        max_y = max(p[1] for p in stroke_pts)
        rx = (max_x - min_x) or 1
        ry = (max_y - min_y) or 1
        pts = [(bx + (p[0] - min_x) / rx * bw, by + (p[1] - min_y) / ry * bh)
               for p in stroke_pts]
        strokes.append({'pts': pts, 'color': nearest_palette(hex_to_rgb(color_hex)),
                        'diag': math.hypot(bw, bh)})

    bodies = [st for st in strokes if st['diag'] >= INK_MARKER_DIAG_FRAC * field_diag]
    markers = [st for st in strokes if st['diag'] < INK_MARKER_DIAG_FRAC * field_diag]

    # attach arrowhead scribbles to the nearest body endpoint
    for mk in markers:
        c = mk['pts'][len(mk['pts']) // 2]
        best = None
        best_d = INK_MARKER_ATTACH_FRAC * field_diag
        for body in bodies:
            for end_i in (0, -1):
                d = dist(c, body['pts'][end_i])
                if d < best_d:
                    best, best_d = (body, end_i), d
        if best:
            body, end_i = best
            body.setdefault('arrow_ends', set()).add(end_i)
        # unattached tiny scribbles are dropped

    for body in bodies:
        pts = body['pts']
        if not chips:
            continue
        k0, d0 = nearest_chip(pts[0])
        k1, d1 = nearest_chip(pts[-1])
        arrow_ends = body.get('arrow_ends', set())
        if d1 < d0:  # flip so the chip-anchored end comes first
            pts = list(reversed(pts))
            arrow_ends = {(-1 if e == 0 else 0) for e in arrow_ends}
            key = k1
        else:
            key = k0
        end = 'ball' if consume_ball(pts[-1]) else 'arrow'
        routes.append({
            'chip': key,
            'color': body['color'],
            'dash': False,
            'end': end,
            'points': [norm(p) for p in downsample(pts)],
        })

    routes, lines = merge_block_caps(routes, lines)

    play = {
        'id': play_meta['import_id'],
        'name': (play_meta['play_name'] or play_meta['play_id'] or 'Untitled')[:60],
        'chips': {k: {'x': norm(c)[0], 'y': norm(c)[1]} for k, c in chips.items()},
        'routes': routes,
        'lines': lines,
        'labels': labels,
        'balls': [{'x': norm(b['c'])[0], 'y': norm(b['c'])[1]}
                  for b in balls if not b['used']],
    }
    return play


# ── main ─────────────────────────────────────────────────────────────────────

def convert(pptx_path, out_path):
    pptx_path = str(pptx_path)
    plays_meta, _sw, _sh = analyze_playbook(pptx_path)
    prs = Presentation(pptx_path)
    pptx_zip = zipfile.ZipFile(pptx_path, 'r')
    theme = load_theme_colors(pptx_zip)

    stem = re.sub(r'[^A-Za-z0-9_-]+', '-', Path(pptx_path).stem)
    doc = {'schema': 1, 'offense': [], 'defense': []}

    for meta in plays_meta:
        if meta['section'] == 'OFFENSE':
            meta['import_id'] = f"import-{stem}-{meta['play_number']:02d}"
            play = convert_play(meta, prs, pptx_zip, theme, 'offense')
            doc['offense'].append(play)
        else:
            meta['import_id'] = f"import-{stem}-D{meta['play_number']}"
            play = convert_play(meta, prs, pptx_zip, theme, 'defense')
            doc['defense'].append(play)

    pptx_zip.close()

    with open(out_path, 'w') as f:
        json.dump(doc, f, indent=1)

    print(f"\nWrote {out_path}")
    print(f"offense={len(doc['offense'])} defense={len(doc['defense'])}")
    for sec in ('offense', 'defense'):
        for i, p in enumerate(doc[sec]):
            print(f"  {sec[:3]} {i+1:>2} {p['name']!r:<22} chips={len(p['chips'])} "
                  f"routes={len(p['routes'])} lines={len(p['lines'])} "
                  f"labels={len(p['labels'])} balls={len(p['balls'])}")
    return doc


def main():
    if len(sys.argv) != 3:
        print('Usage: python pptx_to_editor.py <playbook.pptx> <out.json>')
        sys.exit(1)
    convert(sys.argv[1], sys.argv[2])


if __name__ == '__main__':
    main()
