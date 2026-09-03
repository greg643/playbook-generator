import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "pipeline"))

from playbook_pipeline import analyze_playbook, find_field_rectangle  # noqa: E402


class AnalyzePlaybookTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)
        self.path = Path(self.temp_dir.name) / "deck.pptx"

    @staticmethod
    def add_title(slide, text):
        box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(3), Inches(0.5))
        box.text_frame.text = text

    @staticmethod
    def add_play(
        slide,
        name,
        field_name=None,
        text_box_name=None,
        field_shape=MSO_SHAPE.RECTANGLE,
    ):
        field = slide.shapes.add_shape(
            field_shape,
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(5),
        )
        if field_name is not None:
            field.name = field_name
        box = slide.shapes.add_textbox(Inches(2), Inches(1.05), Inches(3), Inches(0.4))
        box.text_frame.text = name
        if text_box_name is not None:
            box.name = text_box_name

    def test_decorated_headers_simple_plays_and_repeated_sections(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        header = presentation.slides.add_slide(blank)
        self.add_title(header, "OFFENSE")
        for index in range(6):
            header.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(index), Inches(2), Inches(0.2), Inches(0.2),
            )

        play_one = presentation.slides.add_slide(blank)
        self.add_play(play_one, "First")

        repeated_header = presentation.slides.add_slide(blank)
        self.add_title(repeated_header, "OFFENSE")
        play_two = presentation.slides.add_slide(blank)
        self.add_play(play_two, "Second")

        no_field = presentation.slides.add_slide(blank)
        no_field.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(8), Inches(5))

        defense_header = presentation.slides.add_slide(blank)
        self.add_title(defense_header, "DEFENSE")
        defense_play = presentation.slides.add_slide(blank)
        self.add_play(defense_play, "Cover")

        presentation.save(self.path)
        plays, _width, _height = analyze_playbook(self.path)

        self.assertEqual([play["filename"] for play in plays], ["01.png", "02.png", "D1.png"])
        self.assertEqual([play["section"] for play in plays], ["OFFENSE", "OFFENSE", "DEFENSE"])

    def test_free_text_case_insensitive_headers(self):
        # Real decks decorate their separators: Greg's deck says "6v6 OFFENSE".
        # A play slide that merely mentions a section (name "Offense Smash",
        # field rectangle present) must stay a play, not become a header.
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        offense_header = presentation.slides.add_slide(blank)
        self.add_title(offense_header, "6v6 Offense")
        named_play = presentation.slides.add_slide(blank)
        self.add_play(named_play, "Offense Smash")

        defense_header = presentation.slides.add_slide(blank)
        self.add_title(defense_header, "DEFENSE — Zone Looks")
        defense_play = presentation.slides.add_slide(blank)
        self.add_play(defense_play, "Cover 2")

        presentation.save(self.path)
        plays, _width, _height = analyze_playbook(self.path)

        self.assertEqual([play["filename"] for play in plays], ["01.png", "D1.png"])
        self.assertEqual(plays[0]["play_name"], "Offense Smash")

    def test_no_headers_treats_recognizable_plays_as_offense(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        cover = presentation.slides.add_slide(blank)
        self.add_title(cover, "Coach Smith Playbook")
        play_one = presentation.slides.add_slide(blank)
        self.add_play(play_one, "First")
        notes = presentation.slides.add_slide(blank)
        self.add_title(notes, "Install notes")
        play_two = presentation.slides.add_slide(blank)
        self.add_play(play_two, "Second")
        presentation.save(self.path)

        warning_sink = []
        plays, _width, _height = analyze_playbook(
            self.path,
            warning_sink=warning_sink,
        )

        self.assertEqual([play["slide_index"] for play in plays], [1, 3])
        self.assertEqual([play["section"] for play in plays], ["OFFENSE", "OFFENSE"])
        self.assertEqual([play["filename"] for play in plays], ["01.png", "02.png"])
        self.assertEqual(warning_sink, [{
            "code": "skipped_no_field_rectangle",
            "slideCount": 1,
        }])

    def test_defense_only_header_treats_prior_plays_as_offense_and_warns(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        cover = presentation.slides.add_slide(blank)
        self.add_title(cover, "Coach Smith Playbook")
        play_one = presentation.slides.add_slide(blank)
        self.add_play(play_one, "First")
        play_two = presentation.slides.add_slide(blank)
        self.add_play(play_two, "Second")
        defense_header = presentation.slides.add_slide(blank)
        self.add_title(defense_header, "DEFENSE")
        defense_play = presentation.slides.add_slide(blank)
        self.add_play(defense_play, "Cover")
        presentation.save(self.path)

        warning_sink = []
        plays, _width, _height = analyze_playbook(
            self.path,
            warning_sink=warning_sink,
        )

        self.assertEqual(
            [play["section"] for play in plays],
            ["OFFENSE", "OFFENSE", "DEFENSE"],
        )
        self.assertEqual(
            [play["filename"] for play in plays],
            ["01.png", "02.png", "D1.png"],
        )
        self.assertEqual(warning_sink, [{
            "code": "assumed_offense_before_defense",
            "playCount": 2,
        }])

    def test_non_play_slides_before_defense_do_not_trigger_warning(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        cover = presentation.slides.add_slide(blank)
        self.add_title(cover, "Coach Smith Playbook")
        template = presentation.slides.add_slide(blank)
        self.add_play(template, "TEMPLATE")
        defense_header = presentation.slides.add_slide(blank)
        self.add_title(defense_header, "DEFENSE")
        defense_play = presentation.slides.add_slide(blank)
        self.add_play(defense_play, "Cover")
        presentation.save(self.path)

        warning_sink = []
        plays, _width, _height = analyze_playbook(
            self.path,
            warning_sink=warning_sink,
        )

        self.assertEqual([play["section"] for play in plays], ["DEFENSE"])
        self.assertEqual(warning_sink, [])

    def test_offense_header_anywhere_preserves_explicit_section_behavior(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        before_sections = presentation.slides.add_slide(blank)
        self.add_play(before_sections, "Do Not Import")
        defense_header = presentation.slides.add_slide(blank)
        self.add_title(defense_header, "DEFENSE")
        defense_play = presentation.slides.add_slide(blank)
        self.add_play(defense_play, "Cover")
        offense_header = presentation.slides.add_slide(blank)
        self.add_title(offense_header, "OFFENSE")
        offense_play = presentation.slides.add_slide(blank)
        self.add_play(offense_play, "Mesh")
        presentation.save(self.path)

        warning_sink = []
        plays, _width, _height = analyze_playbook(
            self.path,
            warning_sink=warning_sink,
        )

        self.assertEqual([play["slide_index"] for play in plays], [2, 4])
        self.assertEqual([play["section"] for play in plays], ["DEFENSE", "OFFENSE"])
        self.assertEqual(warning_sink, [{
            "code": "skipped_before_first_divider",
            "slideCount": 1,
        }])

    def test_missing_field_inside_explicit_section_produces_warning(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]

        offense_header = presentation.slides.add_slide(blank)
        self.add_title(offense_header, "OFFENSE")
        missing_field = presentation.slides.add_slide(blank)
        self.add_title(missing_field, "Missing Field")
        valid_play = presentation.slides.add_slide(blank)
        self.add_play(valid_play, "Mesh")
        presentation.save(self.path)

        warning_sink = []
        plays, _width, _height = analyze_playbook(
            self.path,
            warning_sink=warning_sink,
        )

        self.assertEqual([play["slide_index"] for play in plays], [2])
        self.assertEqual(warning_sink, [{
            "code": "skipped_no_field_rectangle",
            "slideCount": 1,
        }])

    def test_no_headers_and_no_fields_raises_field_guideline(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        cover = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.add_title(cover, "Coach Smith Playbook")
        presentation.save(self.path)

        with self.assertRaisesRegex(ValueError, "No play slides with a field rectangle"):
            analyze_playbook(self.path)

    def test_more_than_sixteen_offense_plays_are_numbered_without_collisions(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        blank = presentation.slide_layouts[6]
        header = presentation.slides.add_slide(blank)
        self.add_title(header, "OFFENSE")
        for index in range(20):
            play = presentation.slides.add_slide(blank)
            self.add_play(play, f"Play {index + 1}")
        presentation.save(self.path)

        plays, _width, _height = analyze_playbook(self.path)

        self.assertEqual(len(plays), 20)
        self.assertEqual(plays[0]["filename"], "01.png")
        self.assertEqual(plays[-1]["filename"], "20.png")
        self.assertEqual(len({play["filename"] for play in plays}), 20)

    def test_renamed_rectangle_is_detected_by_geometry(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        play = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.add_play(play, "Renamed Field", field_name="Field")
        presentation.save(self.path)

        plays, _width, _height = analyze_playbook(self.path)

        self.assertEqual(len(plays), 1)
        self.assertEqual(plays[0]["section"], "OFFENSE")

    def test_rectangle_family_shapes_are_detected_as_fields(self):
        rectangle_family = (
            MSO_SHAPE.RECTANGLE,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            MSO_SHAPE.ROUND_1_RECTANGLE,
            MSO_SHAPE.ROUND_2_DIAG_RECTANGLE,
            MSO_SHAPE.ROUND_2_SAME_RECTANGLE,
            MSO_SHAPE.SNIP_1_RECTANGLE,
            MSO_SHAPE.SNIP_2_DIAG_RECTANGLE,
            MSO_SHAPE.SNIP_2_SAME_RECTANGLE,
            MSO_SHAPE.SNIP_ROUND_RECTANGLE,
        )

        for field_shape in rectangle_family:
            with self.subTest(field_shape=field_shape):
                presentation = Presentation()
                presentation.slides._sldIdLst.clear()
                play = presentation.slides.add_slide(presentation.slide_layouts[6])
                self.add_play(play, "Compatible Field", field_shape=field_shape)
                presentation.save(self.path)

                plays, _width, _height = analyze_playbook(self.path)

                self.assertEqual(len(plays), 1)
                self.assertEqual(plays[0]["section"], "OFFENSE")

    def test_rectangle_name_falls_back_for_unrecognized_auto_shape(self):
        class FutureRectangleShape:
            shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE
            name = "Rounded Rectangle 3"
            left = 10
            top = 20
            width = 300
            height = 200

            @property
            def auto_shape_type(self):
                raise ValueError("unknown preset geometry")

        self.assertEqual(
            find_field_rectangle([FutureRectangleShape()]),
            (10, 20, 310, 220),
        )

    def test_renamed_text_box_still_supplies_play_name(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        play = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.add_play(play, "Mesh", text_box_name="Play label")
        presentation.save(self.path)

        plays, _width, _height = analyze_playbook(self.path)

        self.assertEqual(plays[0]["play_name"], "Mesh")

    def test_misleading_rectangle_name_on_non_rectangle_is_ignored(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1), Inches(1), Inches(8), Inches(5),
        )
        oval.name = "Rectangle"

        self.assertIsNone(find_field_rectangle(slide.shapes))

    def test_sections_without_plays_raises_guideline(self):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        header = presentation.slides.add_slide(presentation.slide_layouts[6])
        self.add_title(header, "Offense")
        presentation.save(self.path)

        with self.assertRaisesRegex(ValueError, "no play slides with a field rectangle"):
            analyze_playbook(self.path)


if __name__ == "__main__":
    unittest.main()
