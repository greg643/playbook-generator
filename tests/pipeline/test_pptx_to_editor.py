import json
import sys
import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / "pipeline"))

from pptx_to_editor import convert  # noqa: E402


class PptxToEditorFormatTests(unittest.TestCase):
    def setUp(self):
        self.temp_dir = tempfile.TemporaryDirectory()
        self.addCleanup(self.temp_dir.cleanup)
        self.pptx_path = Path(self.temp_dir.name) / "playbook.pptx"
        self.json_path = Path(self.temp_dir.name) / "playbook.json"

    @staticmethod
    def add_section_slide(presentation, section):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(4), Inches(0.6)
        )
        title.text_frame.text = section.upper()

    @staticmethod
    def add_play_slide(presentation, name, player_labels):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1),
            Inches(1),
            Inches(8),
            Inches(5),
        )
        title = slide.shapes.add_textbox(
            Inches(1.25), Inches(1.1), Inches(4), Inches(0.45)
        )
        title.text_frame.text = name

        spacing = 6.5 / max(1, len(player_labels) - 1)
        for index, label in enumerate(player_labels):
            chip = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.45 + index * spacing),
                Inches(3.5),
                Inches(0.55),
                Inches(0.55),
            )
            chip.text_frame.text = label

    def convert_deck(self, sections):
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        for section, plays in sections:
            self.add_section_slide(presentation, section)
            for name, labels in plays:
                self.add_play_slide(presentation, name, labels)
        presentation.save(self.pptx_path)

        returned = convert(self.pptx_path, self.json_path)
        written = json.loads(self.json_path.read_text())
        self.assertEqual(written, returned)
        return written

    def test_exact_five_player_lineups_emit_editable_schema_two(self):
        doc = self.convert_deck([
            ("offense", [("Center Choice", ["1", "2", "C", "3", "QB"])]),
            ("defense", [("Five Across", ["1", "2", "3", "4", "5"])]),
        ])

        self.assertEqual(doc["schema"], 2)
        self.assertEqual(doc["defaultPlayersPerSide"], 5)

        offense = doc["offense"][0]
        self.assertEqual(offense["playersPerSide"], 5)
        self.assertEqual(set(offense["chips"]), {"1", "2", "3", "C", "QB"})
        self.assertNotIn("C", {label["text"] for label in offense["labels"]})

        defense = doc["defense"][0]
        self.assertEqual(defense["playersPerSide"], 5)
        self.assertEqual(set(defense["chips"]), {"1", "2", "3", "4", "5"})

    def test_legacy_five_player_defense_n_is_normalized_to_five(self):
        doc = self.convert_deck([
            ("defense", [("Legacy Five", ["1", "2", "3", "4", "N"])]),
        ])

        defense = doc["defense"][0]
        self.assertEqual(doc["defaultPlayersPerSide"], 5)
        self.assertEqual(defense["playersPerSide"], 5)
        self.assertEqual(set(defense["chips"]), {"1", "2", "3", "4", "5"})
        self.assertNotIn("N", defense["chips"])
        self.assertNotIn("N", {label["text"] for label in defense["labels"]})

    def test_existing_six_player_lineups_remain_six(self):
        doc = self.convert_deck([
            ("offense", [("Six Offense", ["1", "2", "3", "4", "5", "QB"])]),
            ("defense", [("Six Defense", ["1", "2", "3", "4", "5", "N"])]),
        ])

        self.assertEqual(doc["defaultPlayersPerSide"], 6)
        offense = doc["offense"][0]
        defense = doc["defense"][0]
        self.assertEqual(offense["playersPerSide"], 6)
        self.assertEqual(defense["playersPerSide"], 6)
        self.assertEqual(set(offense["chips"]), {"1", "2", "3", "4", "5", "QB"})
        self.assertEqual(set(defense["chips"]), {"1", "2", "3", "4", "5", "N"})

    def test_ambiguous_or_incomplete_lineups_fall_back_to_six(self):
        doc = self.convert_deck([
            ("offense", [("Extra Mark", ["1", "2", "3", "4", "C", "QB"])]),
            ("defense", [("Missing Nose", ["1", "2", "3", "4"])]),
        ])

        offense = doc["offense"][0]
        defense = doc["defense"][0]
        self.assertEqual(offense["playersPerSide"], 6)
        self.assertEqual(defense["playersPerSide"], 6)
        self.assertNotIn("C", offense["chips"])
        self.assertIn("C", {label["text"] for label in offense["labels"]})

    def test_dominant_format_sets_default_and_ties_prefer_six(self):
        dominant_five = self.convert_deck([
            ("offense", [
                ("Five One", ["1", "2", "3", "C", "QB"]),
                ("Five Two", ["1", "2", "3", "C", "QB"]),
                ("Six One", ["1", "2", "3", "4", "5", "QB"]),
            ]),
        ])
        self.assertEqual(dominant_five["defaultPlayersPerSide"], 5)

        tied = self.convert_deck([
            ("offense", [
                ("Five", ["1", "2", "3", "C", "QB"]),
                ("Six", ["1", "2", "3", "4", "5", "QB"]),
            ]),
        ])
        self.assertEqual(tied["defaultPlayersPerSide"], 6)


if __name__ == "__main__":
    unittest.main()
